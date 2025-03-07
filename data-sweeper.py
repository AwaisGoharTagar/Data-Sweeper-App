import streamlit as st
import pandas as pd
import os
from io import BytesIO
import docx
from PyPDF2 import PdfReader
import pptx
from pptx.util import Inches
from fpdf import FPDF

# Set up the App
st.set_page_config(page_title="Data Sweeper", layout="wide")
st.title("Data Sweeper")
st.write("Convert, clean, and visualize your files with built-in processing!")

# File Upload
uploaded_files = st.file_uploader("Upload your files:", type=["csv", "xlsx", "docx", "pptx", "pdf"], accept_multiple_files=True)

if uploaded_files:
    for file in uploaded_files:
        file_ext = os.path.splitext(file.name)[-1].lower()

        # Process Different File Types
        if file_ext == ".csv":
            df = pd.read_csv(file)
            preview_content = df.head()  # Keep dataframe for preview
        elif file_ext == ".xlsx":
            df = pd.read_excel(file, engine="openpyxl")
            preview_content = df.head()
        elif file_ext == ".docx":
            doc = docx.Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])
            df = pd.DataFrame({"Text": text.split("\n")})  # Convert to table
            preview_content = text  # Full text preview
        elif file_ext == ".pdf":
            pdf_reader = PdfReader(file)
            text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
            df = pd.DataFrame({"Text": text.split("\n")})  # Convert to table
            preview_content = text  # Full text preview
        elif file_ext == ".pptx":
            presentation = pptx.Presentation(file)
            slides_text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
            df = pd.DataFrame({"Text": slides_text.split("\n")})  # Convert to table
            preview_content = slides_text  # Full text preview
        else:
            st.error(f"Unsupported file type: {file_ext}")
            continue

        # File Info
        st.write(f"**File Name:** {file.name}")
        st.write(f"**File Size:** {file.size / 1024:.2f} KB")

        # Show Full Preview
        st.subheader("Full File Preview")
        if isinstance(preview_content, str):
            st.text_area(f"Content of {file.name}", preview_content, height=300)  # Full text preview
        else:
            st.dataframe(preview_content)  # DataFrame preview for CSV/Excel

        # **🔹 Data Cleaning Options (Only for CSV & Excel)**
        if file_ext in [".csv", ".xlsx"]:
            st.subheader("Data Cleaning Options")
            if st.checkbox(f"Enable Cleaning for {file.name}"):
                col1, col2 = st.columns(2)
                with col1:
                    if st.button(f"Remove Duplicates from {file.name}"):
                        df.drop_duplicates(inplace=True)
                        st.success("Duplicates Removed!")
                with col2:
                    if st.button(f"Fill Missing Values for {file.name}"):
                        numeric_cols = df.select_dtypes(include=['number']).columns
                        df[numeric_cols] = df[numeric_cols].fillna(df[numeric_cols].mean())
                        st.success("Missing Values Filled!")

            # **🔹 Select Columns to Keep**
            st.subheader("Select Columns to Keep")
            columns = st.multiselect(f"Choose Columns for {file.name}", df.columns, default=df.columns)
            df = df[columns]

            # **🔹 Data Visualization**
            st.subheader("Data Visualization")
            if st.checkbox(f"Show Visualizations for {file.name}"):
                chart_type = st.selectbox("Choose a Chart Type", ["Bar Chart", "Line Chart", "Histogram"], key=file.name)
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) >= 2:
                    if chart_type == "Bar Chart":
                        st.bar_chart(df[numeric_cols])
                    elif chart_type == "Line Chart":
                        st.line_chart(df[numeric_cols])
                    elif chart_type == "Histogram":
                        st.write(df[numeric_cols].hist())
                        st.pyplot()
                else:
                    st.warning("Not enough numeric columns for visualization.")

        # **🔹 File Conversion**
        st.subheader("File Conversion")
        conversion_type = st.radio(f"Convert {file.name} to:", ["CSV", "Excel", "Word", "PowerPoint", "PDF"], key=file.name + "conversion")

        if st.button(f"Convert {file.name}"):
            buffer = BytesIO()

            if conversion_type == "CSV":
                df.to_csv(buffer, index=False)
                file.name = file.name.replace(file_ext, ".csv")
                mime_type = "text/csv"

            elif conversion_type == "Excel":
                df.to_excel(buffer, index=False, engine="openpyxl")
                file.name = file.name.replace(file_ext, ".xlsx")
                mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

            elif conversion_type == "Word":
                doc = docx.Document()
                doc.add_heading("Converted Data", level=1)
                for text in df["Text"]:
                    doc.add_paragraph(text)
                doc.save(buffer)
                file.name = file.name.replace(file_ext, ".docx")
                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

            elif conversion_type == "PowerPoint":
                ppt = pptx.Presentation()
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                text_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(3))
                text_frame = text_box.text_frame
                text_frame.text = "Converted Data"
                for text in df["Text"]:
                    p = text_frame.add_paragraph()
                    p.text = text
                ppt.save(buffer)
                file.name = file.name.replace(file_ext, ".pptx")
                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

            elif conversion_type == "PDF":
                pdf = FPDF()
                pdf.set_auto_page_break(auto=True, margin=15)
                pdf.add_page()
                pdf.set_font("Arial", size=10)
                
                for text in df["Text"]:
                    pdf.multi_cell(190, 10, txt=text)
                
                from io import BytesIO
                buffer = BytesIO()  # Create a BytesIO buffer

                pdf_output = pdf.output(dest="S").encode("latin1")  # Convert PDF to binary
                buffer.write(pdf_output)  # Write to BytesIO
                buffer.seek(0)  # Reset buffer position

                file.name = file.name.replace(file_ext, ".pdf")
                mime_type = "application/pdf"

                # Save locally (Optional)
                pdf.output("output.pdf")

            buffer.seek(0)
            
            # Download Button
            st.download_button(
                label=f"Download {file.name} as {conversion_type}",
                data=buffer,
                file_name=file.name,
                mime=mime_type
            )

st.success("All Files Processed!")
